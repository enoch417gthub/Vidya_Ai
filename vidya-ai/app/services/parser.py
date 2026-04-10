# app/services/parser.py
# ============================================================
# Extracts raw text from uploaded documents
# Supports: PDF, PPTX, DOCX, TXT, and scanned images (via OCR)
# ============================================================
import os
from pathlib import Path
from typing import List, Dict
from loguru import logger
import fitz  # PyMuPDF — fast PDF parser
import pytesseract  # OCR for scanned/image PDFs
from PIL import Image
from pptx import Presentation
from docx import Document


def extract_from_pdf(filepath: str) -> List[Dict]:
    '''
    Extract text from a PDF file page by page.
    Falls back to OCR if page has no extractable text (scanned PDF).
    Returns: list of dicts with 'page', 'text' keys
    '''
    pages = []
    doc = fitz.open(filepath)

    for page_num, page in enumerate(doc):
        text = page.get_text('text')  # Try direct text extraction first

        # If text is too short, it's probably a scanned page — use OCR
        if len(text.strip()) < 50:
            logger.debug(f'Page {page_num+1}: using OCR fallback')
            # Render page as image at 200 DPI for better OCR accuracy
            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(img, lang='eng')  # Add '+hin' for Hindi

        if text.strip():  # Only include pages with actual content
            pages.append({'page': page_num + 1, 'text': text.strip()})

    doc.close()
    return pages


def extract_from_pptx(filepath: str) -> List[Dict]:
    '''
    Extract text from all slides in a PowerPoint file.
    Reads text from text boxes and table cells.
    '''
    prs = Presentation(filepath)
    slides = []

    for slide_num, slide in enumerate(prs.slides):
        slide_texts = []

        for shape in slide.shapes:
            # Text frames (regular text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = ' '.join([run.text for run in para.runs])
                    if text.strip():
                        slide_texts.append(text.strip())

            # Tables inside slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = ' | '.join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        slide_texts.append(row_text)

        if slide_texts:
            slides.append({
                'page': slide_num + 1,
                'text': '\n'.join(slide_texts)
            })

    return slides


def extract_from_docx(filepath: str) -> List[Dict]:
    '''Extract text from a Word document (treated as one 'page')'''
    doc = Document(filepath)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    full_text = '\n'.join(paragraphs)
    return [{'page': 1, 'text': full_text}]


def extract_text(filepath: str) -> List[Dict]:
    '''
    Auto-detect file type and extract text.
    Main entry point — call this for any uploaded document.
    '''
    ext = Path(filepath).suffix.lower()
    logger.info(f'Extracting text from: {os.path.basename(filepath)} ({ext})')

    if ext == '.pdf':
        return extract_from_pdf(filepath)
    elif ext in ['.pptx', '.ppt']:
        return extract_from_pptx(filepath)
    elif ext in ['.docx', '.doc']:
        return extract_from_docx(filepath)
    elif ext == '.txt':
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return [{'page': 1, 'text': f.read()}]
    else:
        raise ValueError(f'Unsupported file type: {ext}')